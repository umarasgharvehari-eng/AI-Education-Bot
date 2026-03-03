import streamlit as st
import re
import faiss
import numpy as np
import pandas as pd
from io import BytesIO
from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from pptx import Presentation
import docx
from groq import Groq

# PDF export (no external font download)
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch


# ---------------------------
# Page Config
# ---------------------------
st.set_page_config(page_title="AI Study Assistant (RAG-based)", layout="wide")


# ---------------------------
# GROQ API Key (Secrets only)
# ---------------------------
if "GROQ_API_KEY" not in st.secrets:
    st.error("GROQ_API_KEY not found in Streamlit Secrets. Add it in App Settings → Secrets.")
    st.stop()

GROQ_API_KEY = st.secrets["GROQ_API_KEY"]


# ---------------------------
# Embedding Model (cached)
# ---------------------------
@st.cache_resource
def load_embedder():
    return SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")


# ---------------------------
# Text extraction
# ---------------------------
def extract_text(uploaded_file):
    name = uploaded_file.name
    ext = name.split(".")[-1].lower()
    raw = uploaded_file.read()
    uploaded_file.seek(0)

    meta = {"source": name}

    if ext == "pdf":
        reader = PdfReader(BytesIO(raw))
        parts = []
        for p in reader.pages:
            parts.append(p.extract_text() or "")
        return "\n".join(parts), meta

    if ext == "pptx":
        prs = Presentation(BytesIO(raw))
        parts = []
        for slide_i, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(f"[Slide {slide_i}] {shape.text}")
        return "\n".join(parts), meta

    if ext == "docx":
        d = docx.Document(BytesIO(raw))
        return "\n".join([p.text for p in d.paragraphs if p.text]), meta

    if ext == "xlsx":
        sheets = pd.read_excel(BytesIO(raw), sheet_name=None)
        parts = []
        for sheet_name, df in sheets.items():
            parts.append(f"[Sheet: {sheet_name}]\n{df.to_csv(index=False)}")
        return "\n\n".join(parts), meta

    if ext in ["txt", "md", "csv"]:
        return raw.decode("utf-8", errors="ignore"), meta

    return "", meta


# ---------------------------
# Chunking
# ---------------------------
def chunk_text(text, chunk_size=3000, overlap=350):
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = min(len(text), start + chunk_size)
        chunks.append(text[start:end])
        if end == len(text):
            break
        start = max(0, end - overlap)
    return chunks


# ---------------------------
# Build KB (FAISS)
# ---------------------------
def build_kb(uploaded_files):
    embedder = load_embedder()

    all_chunks = []
    metadatas = []

    for f in uploaded_files:
        text, meta = extract_text(f)
        for i, ch in enumerate(chunk_text(text)):
            all_chunks.append(ch)
            metadatas.append({**meta, "chunk_no": i})

    if not all_chunks:
        raise ValueError("No readable text found in uploaded files.")

    emb = embedder.encode(all_chunks, normalize_embeddings=True)
    emb = np.array(emb, dtype="float32")

    index = faiss.IndexFlatIP(emb.shape[1])
    index.add(emb)

    return {"index": index, "chunks": all_chunks, "metadatas": metadatas}


def retrieve(question, kb, k=8):
    embedder = load_embedder()
    q = embedder.encode([question], normalize_embeddings=True)
    q = np.array(q, dtype="float32")
    D, I = kb["index"].search(q, k)

    hits = []
    for idx in I[0]:
        if idx == -1:
            continue
        hits.append((kb["chunks"][idx], kb["metadatas"][idx]))
    return hits


# ---------------------------
# GROQ Call
# ---------------------------
def call_groq(system_prompt, user_prompt, model="llama-3.3-70b-versatile"):
    client = Groq(api_key=GROQ_API_KEY)
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_prompt.strip()},
            {"role": "user", "content": user_prompt.strip()},
        ],
        temperature=0.2,
    )
    return resp.choices[0].message.content


# ---------------------------
# PDF Export (ReportLab) + Sanitizer
# ---------------------------
def sanitize_for_pdf(text: str) -> str:
    """
    ReportLab default fonts are not fully unicode-friendly.
    We sanitize to ASCII-safe content to prevent crashes.
    """
    replacements = {
        "→": "->",
        "⇒": "=>",
        "•": "-",
        "−": "-",
        "×": "x",
        "°": " deg ",
        "²": "^2",
        "³": "^3",
        "₄": "4",
        "₃": "3",
        "₁": "1",
        "₀": "0",
    }
    for k, v in replacements.items():
        text = text.replace(k, v)

    # Remove remaining non-ascii
    text = re.sub(r"[^\x00-\x7F]+", "", text)
    return text


def make_pdf_bytes(text: str) -> bytes:
    """
    Creates a PDF in memory using ReportLab (no internet calls).
    """
    text = sanitize_for_pdf(text)

    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter

    left_margin = 0.75 * inch
    top_margin = 0.75 * inch
    line_height = 14
    y = height - top_margin

    c.setFont("Helvetica", 11)

    # Basic word wrap
    max_width = width - 2 * left_margin

    def draw_wrapped_line(line, y_pos):
        words = line.split(" ")
        current = ""
        for w in words:
            test = (current + " " + w).strip()
            if c.stringWidth(test, "Helvetica", 11) <= max_width:
                current = test
            else:
                c.drawString(left_margin, y_pos, current)
                y_pos -= line_height
                current = w
                if y_pos < top_margin:
                    c.showPage()
                    c.setFont("Helvetica", 11)
                    y_pos = height - top_margin
        if current:
            c.drawString(left_margin, y_pos, current)
            y_pos -= line_height
        return y_pos

    for line in text.splitlines():
        if not line.strip():
            y -= line_height
        else:
            y = draw_wrapped_line(line, y)

        if y < top_margin:
            c.showPage()
            c.setFont("Helvetica", 11)
            y = height - top_margin

    c.save()
    buffer.seek(0)
    return buffer.read()


# ---------------------------
# Session State
# ---------------------------
if "messages" not in st.session_state:
    st.session_state.messages = []  # [{"role":"user/assistant","content":"..."}]
if "kb" not in st.session_state:
    st.session_state.kb = None


# ---------------------------
# Sidebar (All controls here)
# ---------------------------
with st.sidebar:
    st.header("Settings")

    education_level = st.selectbox(
        "Select your education level",
        ["Primary School", "Middle School", "Secondary School", "College/High School", "Undergraduate", "Graduate"],
        index=0,
    )

    st.markdown("---")
    st.subheader("Upload study materials")
    uploaded_files = st.file_uploader(
        "Supported: PDF, PPTX, DOCX, XLSX, TXT, MD, CSV",
        accept_multiple_files=True,
        type=["pdf", "pptx", "docx", "xlsx", "txt", "md", "csv"],
    )

    colA, colB = st.columns(2)
    with colA:
        build_btn = st.button("Build KB", use_container_width=True)
    with colB:
        clear_kb_btn = st.button("Clear KB", use_container_width=True)

    colC, colD = st.columns(2)
    with colC:
        clear_chat_btn = st.button("Clear Chat", use_container_width=True)
    with colD:
        pass

    if build_btn:
        if not uploaded_files:
            st.warning("Upload files first.")
        else:
            with st.spinner("Building knowledge base (extracting + embeddings + FAISS)..."):
                try:
                    st.session_state.kb = build_kb(uploaded_files)
                    st.success("KB is ready.")
                except Exception as e:
                    st.session_state.kb = None
                    st.error(f"KB build failed: {e}")

    if clear_kb_btn:
        st.session_state.kb = None
        st.success("KB cleared.")

    if clear_chat_btn:
        st.session_state.messages = []
        st.success("Chat cleared.")

    st.markdown("---")
    if st.session_state.kb is None:
        st.info("KB status: Not built yet")
    else:
        st.success("KB status: Ready")


# ---------------------------
# Main UI (Chat only)
# ---------------------------
st.title("AI Study Assistant (RAG-based)")
st.caption("Chat with your study materials (English only).")

# Show chat history
for msg in st.session_state.messages:
    with st.chat_message(msg["role"]):
        st.write(msg["content"])

# Chat input
user_q = st.chat_input("Type your question...")

if user_q:
    # Store user message
    st.session_state.messages.append({"role": "user", "content": user_q})
    with st.chat_message("user"):
        st.write(user_q)

    # If KB not ready
    if st.session_state.kb is None:
        assistant_text = "Knowledge base is not ready. Please upload files and click **Build KB** in the sidebar."
        st.session_state.messages.append({"role": "assistant", "content": assistant_text})
        with st.chat_message("assistant"):
            st.write(assistant_text)
    else:
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                hits = retrieve(user_q, st.session_state.kb, k=10)

                context_blocks = []
                cite_files = []
                for ch, meta in hits:
                    context_blocks.append(f"[Source: {meta['source']}]\n{ch}")
                    cite_files.append(meta["source"])
                cite_files = sorted(set(cite_files))

                # ENGLISH ONLY system prompt
                system_prompt = f"""
You are an AI Study Assistant.

STRICT REQUIREMENTS:
- Respond ONLY in English.
- Answer ONLY using the provided CONTEXT (no outside knowledge).
- If the answer is not in the context, say: "I don't have enough information in the uploaded materials."
- Explain in a way appropriate for: {education_level}.
- Structure the answer with clear steps when needed.

OUTPUT FORMAT:
1) Answer
2) Short Summary (6-10 lines)
3) Sources: list filenames only
"""

                user_prompt = f"""
QUESTION:
{user_q}

CONTEXT:
{'\n\n---\n\n'.join(context_blocks)}
"""

                answer = call_groq(system_prompt, user_prompt)

                # Ensure sources line exists
                if "Sources:" not in answer and cite_files:
                    answer = answer.strip() + "\n\nSources: " + ", ".join(cite_files)

            st.write(answer)

            # Downloads (English-only output)
            st.download_button("Download Markdown", answer, file_name="summary.md")

            pdf_bytes = make_pdf_bytes(answer)
            st.download_button("Download PDF", pdf_bytes, file_name="summary.pdf", mime="application/pdf")

        st.session_state.messages.append({"role": "assistant", "content": answer})
